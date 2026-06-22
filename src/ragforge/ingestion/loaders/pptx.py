from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ragforge.ingestion.loaders.base import BaseDocumentLoader, register_loader
from ragforge.ingestion.models import DocumentUnit, LoadedBlock
from ragforge.ingestion.processors import BaseUnitProcessor
from ragforge.ingestion.utils import hash_file, normalize_whitespace


@register_loader(".pptx", ".ppt")
class PptxLoader(BaseUnitProcessor, BaseDocumentLoader):
    def _build_units(self, file_path: str) -> list[DocumentUnit]:
        prs = Presentation(file_path)
        units: list[DocumentUnit] = []
        file_hash = hash_file(file_path)

        for idx, slide in enumerate(prs.slides):
            slide_title = f"Slide {idx + 1}"
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()

            text_parts = []
            image_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape == slide.shapes.title:
                        continue
                    text_parts.append(shape.text.strip())
                if getattr(shape, "shape_type", None) == 13:
                    image_count += 1

            body_text = normalize_whitespace("\n".join(text_parts))
            text = normalize_whitespace(f"Title: {slide_title}\n{body_text}".strip())
            if not text:
                continue

            units.append(
                DocumentUnit(
                    doc_id=file_hash,
                    source_path=file_path,
                    source_name=Path(file_path).name,
                    file_type="pptx",
                    unit_num=idx + 1,
                    raw_text=text,
                    image_count=image_count,
                    metadata={
                        "source": file_path,
                        "filename": Path(file_path).name,
                        "file_type": "pptx",
                        "slide_number": idx + 1,
                        "slide_title": slide_title,
                        "image_count": image_count,
                    },
                    unit_kind="slide",
                )
            )

        return units

    def load(self, file_path: str, config_path: str | None = None) -> list[LoadedBlock]:
        units = self._build_units(file_path)
        return self.process_units(units, config_path=config_path)

    async def load_async(
        self, file_path: str, config_path: str | None = None
    ) -> list[LoadedBlock]:
        units = self._build_units(file_path)
        return await self.process_units_async(units, config_path=config_path)
