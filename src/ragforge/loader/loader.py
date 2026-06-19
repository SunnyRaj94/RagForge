import os
import fitz  # PyMuPDF
import openpyxl
from pptx import Presentation
import easyocr
import yaml
from typing import List, Dict, Any
from src.ragforge.utils.logging_hooks import observe_stage

# Lazy EasyOCR reader initializer
_ocr_reader = None


def get_ocr_reader(languages: List[str] = None):
    global _ocr_reader
    if _ocr_reader is None:
        if languages is None:
            languages = ["en"]
        _ocr_reader = easyocr.Reader(languages, gpu=True)
    return _ocr_reader


def load_chunking_config(config_path: str = None) -> Dict[str, Any]:
    if config_path is None:
        from src.ragforge.config import get_config_path

        config_path = get_config_path()

    if not os.path.exists(config_path):
        return {}
    with open(config_path, "r") as f:
        return yaml.safe_load(f) or {}


def table_to_markdown(cells: List[List[Any]]) -> str:
    if not cells or not any(cells):
        return ""
    # Find max number of columns
    max_cols = max(len(row) for row in cells if row)
    markdown_lines = []

    # Process header
    header = cells[0] if len(cells) > 0 else []
    header_str = "| " + " | ".join(
        str(cell or "").replace("\n", " ") for cell in header
    )
    # Pad header if shorter than max columns
    if len(header) < max_cols:
        header_str += " | " + " | ".join("" for _ in range(max_cols - len(header)))
    header_str += " |"
    markdown_lines.append(header_str)

    # Process separator
    markdown_lines.append("| " + " | ".join("---" for _ in range(max_cols)) + " |")

    # Process data rows
    for row in cells[1:]:
        if not row:
            continue
        row_str = "| " + " | ".join(str(cell or "").replace("\n", " ") for cell in row)
        if len(row) < max_cols:
            row_str += " | " + " | ".join("" for _ in range(max_cols - len(row)))
        row_str += " |"
        markdown_lines.append(row_str)

    return "\n".join(markdown_lines)


@observe_stage("parse_pdf")
def parse_pdf(file_path: str, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Parses PDF using PyMuPDF (fitz).
    - Digital text extraction.
    - Fallback to OCR if digital text is empty.
    - Extract tables as markdown grids if configured.
    """
    pdf_config = config.get("pdf", {})
    extract_tables = pdf_config.get("extract_tables", True)
    ocr_fallback = pdf_config.get("ocr_fallback", True)
    ocr_languages = pdf_config.get("ocr_language", ["en"])

    doc = fitz.open(file_path)
    chunks = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        # 1. Digital Text Extraction
        digital_text = page.get_text().strip()

        # 2. Table Extraction
        table_text = ""
        if extract_tables:
            try:
                tables = page.find_tables()
                for table in tables:
                    cells = table.extract()
                    if cells:
                        table_text += "\n\n" + table_to_markdown(cells) + "\n\n"
            except Exception:
                pass  # Fallback if table extraction fails

        # 3. Decision for OCR
        text = digital_text
        ocr_run = False
        if not text and ocr_fallback:
            # Render page to image bytes for OCR
            pix = page.get_pixmap(dpi=150)
            img_path = f"{file_path}_page_{page_num}.png"
            pix.save(img_path)

            try:
                reader = get_ocr_reader(ocr_languages)
                ocr_result = reader.readtext(img_path)
                text = " ".join([res[1] for res in ocr_result]).strip()
                ocr_run = True
            finally:
                if os.path.exists(img_path):
                    os.remove(img_path)

        full_content = text
        if table_text:
            full_content += "\n" + table_text.strip()

        if full_content.strip():
            chunks.append(
                {
                    "text": full_content.strip(),
                    "metadata": {
                        "source": file_path,
                        "filename": os.path.basename(file_path),
                        "file_type": "pdf",
                        "page": page_num + 1,
                        "ocr_run": ocr_run,
                    },
                }
            )

    return chunks


@observe_stage("parse_excel")
def parse_excel(file_path: str, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Parses Excel using openpyxl.
    - Groups rows into batch sizes (row_batch_size)
    - Prepends headers to each batch chunk.
    """
    excel_config = config.get("excel", {})
    row_batch_size = excel_config.get("row_batch_size", 15)
    header_row_idx = excel_config.get("header_row", 1) - 1  # Convert to 0-indexed
    sheet_mode = excel_config.get("sheet_mode", "all")

    wb = openpyxl.load_workbook(file_path, data_only=True)
    chunks = []

    sheets_to_parse = wb.sheetnames
    if sheet_mode == "first" and sheets_to_parse:
        sheets_to_parse = [sheets_to_parse[0]]

    for sheet_name in sheets_to_parse:
        sheet = wb[sheet_name]
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue

        # Extract headers from the designated header row
        headers = []
        if header_row_idx < len(rows):
            headers = [
                str(cell) if cell is not None else "" for cell in rows[header_row_idx]
            ]

        # Data rows start after header row
        data_rows = rows[header_row_idx + 1 :]
        if not data_rows:
            continue

        # Batch data rows
        for i in range(0, len(data_rows), row_batch_size):
            batch = data_rows[i : i + row_batch_size]

            # Format batch as markdown representation
            batch_lines = []
            if headers:
                batch_lines.append("| " + " | ".join(headers) + " |")
                batch_lines.append("| " + " | ".join("---" for _ in headers) + " |")

            for r in batch:
                row_cells = [str(cell) if cell is not None else "" for cell in r]
                # Pad/trim to headers length
                if len(row_cells) < len(headers):
                    row_cells += [""] * (len(headers) - len(row_cells))
                elif len(row_cells) > len(headers) and headers:
                    row_cells = row_cells[: len(headers)]
                batch_lines.append("| " + " | ".join(row_cells) + " |")

            chunk_text = f"Sheet: {sheet_name}\n" + "\n".join(batch_lines)

            chunks.append(
                {
                    "text": chunk_text,
                    "metadata": {
                        "source": file_path,
                        "filename": os.path.basename(file_path),
                        "file_type": "excel",
                        "sheet": sheet_name,
                        "start_row": i + header_row_idx + 2,
                        "end_row": i + header_row_idx + 1 + len(batch),
                    },
                }
            )

    return chunks


@observe_stage("parse_pptx")
def parse_pptx(file_path: str, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Parses PowerPoint presentation using python-pptx.
    - Treat each slide as a self-contained chunk.
    - Prepend slide titles.
    """
    pptx_config = config.get("pptx", {})
    slide_as_chunk = pptx_config.get("slide_as_chunk", True)
    prepend_slide_title = pptx_config.get("prepend_slide_title", True)

    prs = Presentation(file_path)
    chunks = []

    for idx, slide in enumerate(prs.slides):
        slide_title = f"Slide {idx + 1}"
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip()

        # Gather text from shapes
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Avoid capturing the title shape text twice
                if shape == slide.shapes.title:
                    continue
                text_parts.append(shape.text.strip())

        slide_content = "\n".join(text_parts).strip()

        full_text = ""
        if prepend_slide_title:
            full_text += f"Title: {slide_title}\n"
        full_text += slide_content

        if full_text.strip():
            chunks.append(
                {
                    "text": full_text.strip(),
                    "metadata": {
                        "source": file_path,
                        "filename": os.path.basename(file_path),
                        "file_type": "pptx",
                        "slide_number": idx + 1,
                        "slide_title": slide_title,
                    },
                }
            )

    return chunks


@observe_stage("parse_image")
def parse_image(file_path: str, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Parses image file using EasyOCR.
    """
    img_config = config.get("image", {})
    ocr_languages = img_config.get("ocr_languages", ["en"])

    reader = get_ocr_reader(ocr_languages)
    result = reader.readtext(file_path)
    text = " ".join([res[1] for res in result]).strip()

    if text:
        return [
            {
                "text": text,
                "metadata": {
                    "source": file_path,
                    "filename": os.path.basename(file_path),
                    "file_type": "image",
                },
            }
        ]
    return []


def load_file(file_path: str, config_path: str = None) -> List[Dict[str, Any]]:
    """
    General entry point for loading any supported document type.
    """
    config = load_chunking_config(config_path)
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path, config)
    elif ext in [".xlsx", ".xls"]:
        return parse_excel(file_path, config)
    elif ext in [".pptx", ".ppt"]:
        return parse_pptx(file_path, config)
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
        return parse_image(file_path, config)
    else:
        # Fallback to plain text reader
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read().strip()
            if content:
                return [
                    {
                        "text": content,
                        "metadata": {
                            "source": file_path,
                            "filename": os.path.basename(file_path),
                            "file_type": "text",
                        },
                    }
                ]
        except Exception:
            pass
        return []
