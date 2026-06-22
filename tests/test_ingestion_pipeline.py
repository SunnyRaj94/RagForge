import os
import pytest
import openpyxl
from pptx import Presentation
from src.ragforge.loader.loader import load_file
from src.ragforge.chunking.chunker import chunk_documents
from src.ragforge.embeddings.embeddings import get_embedding
from src.ragforge.index.indexer import (
    create_collection,
    upsert_documents,
    get_qdrant_client,
)


def test_excel_parsing():
    excel_path = "test_mock.xlsx"
    # Create simple Excel file
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["Name", "Role", "Description"])
    ws.append(["Alice", "Developer", "Writes clean python code."])
    ws.append(["Bob", "Manager", "Coordinates tasks and schedules."])
    wb.save(excel_path)

    try:
        docs = load_file(excel_path)
        assert len(docs) == 1
        assert "Alice" in docs[0]["text"]
        assert "Developer" in docs[0]["text"]
        assert "TestSheet" in docs[0]["text"]
        assert docs[0]["metadata"]["file_type"] == "excel"
    finally:
        if os.path.exists(excel_path):
            os.remove(excel_path)


def test_pptx_parsing():
    pptx_path = "test_mock.pptx"
    # Create simple PowerPoint
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Ingestion Title"
    txBox = slide.shapes.add_textbox(100, 100, 100, 100)
    tf = txBox.text_frame
    tf.text = "This is body content for slide 1."
    prs.save(pptx_path)

    try:
        docs = load_file(pptx_path)
        assert len(docs) == 1
        assert "Ingestion Title" in docs[0]["text"]
        assert "body content" in docs[0]["text"]
        assert docs[0]["metadata"]["file_type"] == "pptx"
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)


def test_end_to_end_pipeline():
    # 1. Create a mock text file
    txt_path = "test_mock.txt"
    with open(txt_path, "w") as f:
        f.write(
            "RagForge is a modular data ingestion framework designed to feed LLMs clean data.\n"
            * 20
        )

    collection_name = "test_pipeline_collection"
    q_client = get_qdrant_client()

    try:
        # Clean collection
        if q_client.collection_exists(collection_name):
            q_client.delete_collection(collection_name)

        # 2. Load
        raw_docs = load_file(txt_path)
        assert len(raw_docs) == 1

        # 3. Chunk
        chunked_docs = chunk_documents(raw_docs)
        assert len(chunked_docs) > 1

        # 4. Embed & Index
        create_res = create_collection(collection_name, 768)
        assert "Collection" in create_res

        upsert_res = upsert_documents(collection_name, chunked_docs)
        assert "Successfully upserted" in upsert_res

        # Verify query in Qdrant
        search_res = q_client.query_points(
            collection_name, query=get_embedding("LLMs clean data"), using="text_vector", limit=1
        )
        assert len(search_res.points) == 1
        assert "RagForge" in search_res.points[0].payload["text"]

    finally:
        if os.path.exists(txt_path):
            os.remove(txt_path)
        if q_client.collection_exists(collection_name):
            q_client.delete_collection(collection_name)


if __name__ == "__main__":
    import sys

    test_excel_parsing()
    test_pptx_parsing()
    test_end_to_end_pipeline()
    print("All ingestion pipeline tests passed successfully!")
